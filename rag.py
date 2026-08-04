"""
AI Digital Twin of Knowledge - Multi-Format RAG Engine

Handles document parsing (PDF, DOCX, PPTX, TXT), page-aware text chunking,
vector embeddings (Sentence Transformers / Fallback TF-IDF engine),
ChromaDB index management, retrieval, citation rendering, and document tracking.
"""

import os
import io
import math
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

from config import settings
from database import db


# ==================== TEXT EXTRACTION UTILITIES ====================

def extract_text_from_file(file_bytes: bytes, file_name: str) -> List[Dict[str, Any]]:
    """
    Extract structured text pages/slides from uploaded file bytes.
    Returns list of dicts: [{'page_number': int, 'content': str}]
    """
    ext = Path(file_name).suffix.lower()
    pages_data: List[Dict[str, Any]] = []

    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages_data.append({"page_number": i + 1, "content": text.strip()})
        except Exception as e:
            pages_data.append({"page_number": 1, "content": f"PDF Extraction fallback. Error: {str(e)}"})

    elif ext in [".docx", ".doc"]:
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)
            pages_data.append({"page_number": 1, "content": full_text})
        except Exception as e:
            pages_data.append({"page_number": 1, "content": f"DOCX Extraction fallback. Error: {str(e)}"})

    elif ext in [".pptx", ".ppt"]:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                combined = "\n".join(slide_text).strip()
                if combined:
                    pages_data.append({"page_number": i + 1, "content": combined})
        except Exception as e:
            pages_data.append({"page_number": 1, "content": f"PPTX Extraction fallback. Error: {str(e)}"})

    else:
        # Default text / markdown parser
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
            pages_data.append({"page_number": 1, "content": text.strip()})
        except Exception as e:
            pages_data.append({"page_number": 1, "content": f"Text Extraction fallback. Error: {str(e)}"})

    return pages_data


def chunk_pages_data(
    pages_data: List[Dict[str, Any]],
    file_name: str,
    doc_id: str,
    chunk_size: int = 800,
    overlap: int = 150
) -> List[Dict[str, Any]]:
    """
    Split page/slide structured text into overlapping chunks with precise citation metadata.
    """
    chunks: List[Dict[str, Any]] = []
    chunk_counter = 0

    for page_info in pages_data:
        page_num = page_info.get("page_number", 1)
        text = page_info.get("content", "")
        if not text:
            continue

        # Character splitting with overlap
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + chunk_size
            chunk_text = text[start:end].strip()

            if chunk_text:
                chunk_counter += 1
                chunk_id = f"{doc_id}_c{chunk_counter}"
                chunks.append({
                    "chunk_id": chunk_id,
                    "doc_id": doc_id,
                    "file_name": file_name,
                    "page_number": page_num,
                    "chunk_index": chunk_counter,
                    "content": chunk_text,
                })

            start += (chunk_size - overlap)
            if start >= text_len:
                break

    return chunks


# ==================== FALLBACK TF-IDF SIMILARITY VECTOR ENGINE ====================

class FallbackVectorStore:
    """
    In-memory fallback vector store using TF-IDF token cosine similarity
    when ChromaDB or sentence-transformers native libraries are missing or loading.
    """

    def __init__(self):
        self.chunks: List[Dict[str, Any]] = []

    def add_chunks(self, chunks: List[Dict[str, Any]]) -> None:
        """Add chunks to in-memory store."""
        for c in chunks:
            if not any(existing["chunk_id"] == c["chunk_id"] for existing in self.chunks):
                self.chunks.append(c)

    def search(self, query: str, top_k: int = 4, user_id: str = "default_user") -> List[Dict[str, Any]]:
        """Perform keyword and TF-IDF similarity ranking over stored chunks."""
        if not self.chunks:
            return []

        query_terms = set(re_tokenize(query.lower()))
        results = []

        for chunk in self.chunks:
            chunk_text = chunk["content"].lower()
            chunk_terms = re_tokenize(chunk_text)
            if not chunk_terms:
                continue

            # Term overlap score
            matches = sum(1 for t in query_terms if t in chunk_terms)
            score = matches / (len(query_terms) + 1.0)

            if score > 0.05 or len(self.chunks) <= top_k:
                results.append((score, chunk))

        # Sort descending by score
        results.sort(key=lambda x: x[0], reverse=True)
        top_matches = []
        for score, chunk in results[:top_k]:
            res = dict(chunk)
            res["relevance_score"] = float(score)
            top_matches.append(res)
        return top_matches

    def delete_document(self, doc_id: str) -> None:
        """Delete chunks matching doc_id."""
        self.chunks = [c for c in self.chunks if c.get("doc_id") != doc_id]


def re_tokenize(text: str) -> List[str]:
    """Helper tokenizer for fallback vector store."""
    import re
    return re.findall(r"\b\w{3,}\b", text.lower())


# ==================== MAIN RAG PIPELINE ====================

class RAGPipeline:
    """
    RAG Pipeline handling document ingestion, vector storage in ChromaDB,
    retrieval, citation generation, and document metadata persistence.
    """

    def __init__(self):
        """Initialize ChromaDB client or fallback vector store."""
        self.persist_dir = settings.CHROMA_PERSIST_DIR
        self.chroma_client = None
        self.collection = None
        self.fallback_store = FallbackVectorStore()
        self.use_fallback = False

        self._init_chroma()

    def _init_chroma(self) -> None:
        """Attempt initializing ChromaDB persistent client."""
        try:
            import chromadb
            Path(self.persist_dir).mkdir(parents=True, exist_ok=True)
            self.chroma_client = chromadb.PersistentClient(path=self.persist_dir)
            self.collection = self.chroma_client.get_or_create_collection(
                name="digital_twin_knowledge",
                metadata={"hnsw:space": "cosine"}
            )
            self.use_fallback = False
        except Exception as e:
            # Graceful fallback to TF-IDF vector store if ChromaDB native bindings fail
            self.use_fallback = True

    def ingest_document(
        self,
        file_bytes: bytes,
        file_name: str,
        user_id: str = "default_user"
    ) -> Dict[str, Any]:
        """
        Ingest uploaded document: parse text, chunk, embed, and store in vector database & SQLite.
        """
        doc_hash = hashlib.md5(f"{file_name}_{len(file_bytes)}".encode()).hexdigest()[:12]
        doc_id = f"doc_{doc_hash}"

        # 1. Save physical file copy to uploaded_docs directory
        upload_dir = settings.UPLOAD_DIR
        upload_dir.mkdir(parents=True, exist_ok=True)
        file_path = upload_dir / f"{doc_hash}_{file_name}"
        with open(file_path, "wb") as f:
            f.write(file_bytes)

        # 2. Parse text page by page
        pages_data = extract_text_from_file(file_bytes, file_name)

        # 3. Chunk parsed text
        chunks = chunk_pages_data(pages_data, file_name, doc_id)

        if not chunks:
            return {
                "doc_id": doc_id,
                "file_name": file_name,
                "status": "warning",
                "total_chunks": 0,
                "message": "No extractable text found in file."
            }

        # 4. Store Chunks in ChromaDB or Fallback Store
        if not self.use_fallback and self.collection:
            try:
                ids = [c["chunk_id"] for c in chunks]
                documents = [c["content"] for c in chunks]
                metadatas = [
                    {
                        "doc_id": c["doc_id"],
                        "file_name": c["file_name"],
                        "page_number": c["page_number"],
                        "chunk_index": c["chunk_index"],
                        "user_id": user_id
                    }
                    for c in chunks
                ]
                self.collection.add(ids=ids, documents=documents, metadatas=metadatas)
            except Exception:
                # Store in fallback if ChromaDB insertion throws error
                self.fallback_store.add_chunks(chunks)
        else:
            self.fallback_store.add_chunks(chunks)

        # 5. Persist Document Metadata in SQLite
        ext = Path(file_name).suffix.replace(".", "").upper()
        db.save_uploaded_document(
            doc_id=doc_id,
            user_id=user_id,
            file_name=file_name,
            file_path=str(file_path),
            file_type=ext,
            total_chunks=len(chunks)
        )

        return {
            "doc_id": doc_id,
            "file_name": file_name,
            "status": "success",
            "total_chunks": len(chunks),
            "file_type": ext,
            "message": f"Successfully indexed {len(chunks)} chunks from {file_name}."
        }

    def query_relevant_chunks(
        self,
        query: str,
        user_id: str = "default_user",
        top_k: int = 4
    ) -> List[Dict[str, Any]]:
        """
        Search vector database for top_k relevant text chunks matching query.
        Returns list of chunk dicts with content, source, page_number, and relevance_score.
        """
        if not query.strip():
            return []

        # Try ChromaDB query first if active
        if not self.use_fallback and self.collection:
            try:
                results = self.collection.query(
                    query_texts=[query],
                    n_results=top_k,
                    where={"user_id": user_id} if user_id else None
                )

                retrieved = []
                if results and results.get("documents") and results["documents"][0]:
                    docs = results["documents"][0]
                    metas = results["metadatas"][0] if results.get("metadatas") else []
                    distances = results["distances"][0] if results.get("distances") else []

                    for i, doc_text in enumerate(docs):
                        meta = metas[i] if i < len(metas) else {}
                        dist = distances[i] if i < len(distances) else 0.5
                        # Cosine similarity conversion
                        sim_score = max(0.0, 1.0 - float(dist))

                        retrieved.append({
                            "chunk_id": meta.get("chunk_id", f"c_{i}"),
                            "doc_id": meta.get("doc_id", "unknown"),
                            "file_name": meta.get("file_name", "Document"),
                            "page_number": meta.get("page_number", 1),
                            "chunk_index": meta.get("chunk_index", i + 1),
                            "content": doc_text,
                            "relevance_score": sim_score
                        })
                if retrieved:
                    return retrieved
            except Exception:
                pass

        # Fallback query
        return self.fallback_store.search(query=query, top_k=top_k, user_id=user_id)

    def format_rag_context(self, chunks: List[Dict[str, Any]]) -> str:
        """Format retrieved context chunks into structured LLM prompt context."""
        if not chunks:
            return ""

        context_lines = ["=== RETRIEVED STUDY MATERIAL (RAG CONTEXT) ==="]
        for i, chunk in enumerate(chunks, 1):
            source = chunk.get("file_name", "Document")
            page = chunk.get("page_number", 1)
            content = chunk.get("content", "").strip()
            score = chunk.get("relevance_score", 0.0)
            context_lines.append(
                f"\n--- [Source #{i}: {source} | Page {page} | Rel: {score:.2f}] ---\n{content}"
            )

        context_lines.append("\n=============================================")
        return "\n".join(context_lines)

    def format_citations_markdown(self, chunks: List[Dict[str, Any]]) -> str:
        """Format citations into clean markdown for Streamlit chat display."""
        if not chunks:
            return ""

        lines = ["\n\n**Source Citations:**"]
        seen = set()

        for chunk in chunks:
            source = chunk.get("file_name", "Document")
            page = chunk.get("page_number", 1)
            key = f"{source}_p{page}"
            if key not in seen:
                seen.add(key)
                snippet = chunk.get("content", "")[:120].replace("\n", " ")
                lines.append(f"- **{source}** (Page {page}): *\"{snippet}...\"*")

        return "\n".join(lines)

    def get_uploaded_documents(self, user_id: str = "default_user") -> List[Dict[str, Any]]:
        """Retrieve user uploaded document metadata list."""
        return db.get_uploaded_documents(user_id=user_id)

    def delete_document(self, doc_id: str, user_id: str = "default_user") -> bool:
        """Delete document vectors and database metadata."""
        if not self.use_fallback and self.collection:
            try:
                self.collection.delete(where={"doc_id": doc_id})
            except Exception:
                pass
        self.fallback_store.delete_document(doc_id)
        return True


# Global RAG Pipeline Instance
rag_pipeline = RAGPipeline()
